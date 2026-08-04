from __future__ import annotations

import csv
import hashlib
import importlib.util
import json
import mimetypes
import os
import uuid
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Iterable

from .multimodal_enrichment import describe_image, enrichment_capabilities, run_ocr


@dataclass
class DocumentBlock:
    block_id: str
    document_id: str
    source: str
    block_type: str
    content: str
    page: int | None = None
    bbox: list[float] | None = None
    section_title: str | None = None
    metadata: dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


def _file_hash(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        while chunk := stream.read(1024 * 1024):
            digest.update(chunk)
    return digest.hexdigest()


def _block(
    path: Path,
    document_id: str,
    block_type: str,
    content: str,
    index: int,
    *,
    page: int | None = None,
    bbox: Iterable[float] | None = None,
    section_title: str | None = None,
    metadata: dict[str, Any] | None = None,
) -> DocumentBlock:
    identity = f"{document_id}:{block_type}:{page}:{index}:{content}"
    block_id = hashlib.sha256(identity.encode("utf-8")).hexdigest()[:24]
    return DocumentBlock(
        block_id=block_id,
        document_id=document_id,
        source=path.name,
        block_type=block_type,
        content=content.strip(),
        page=page,
        bbox=[round(float(value), 2) for value in bbox] if bbox else None,
        section_title=section_title,
        metadata=metadata or {},
    )


def _markdown_table(rows: list[list[Any]]) -> str:
    normalized = [
        [str(cell or "").replace("|", "\\|").replace("\n", " ").strip() for cell in row]
        for row in rows
        if row
    ]
    if not normalized:
        return ""
    width = max(len(row) for row in normalized)
    normalized = [row + [""] * (width - len(row)) for row in normalized]
    header = normalized[0]
    body = normalized[1:]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def _parse_text(path: Path, document_id: str) -> list[DocumentBlock]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    blocks: list[DocumentBlock] = []
    section: str | None = None
    for index, paragraph in enumerate(part.strip() for part in text.split("\n\n")):
        if not paragraph:
            continue
        first = next((line.strip() for line in paragraph.splitlines() if line.strip()), "")
        is_title = first.startswith("#") or (
            len(first) <= 80 and first.startswith("第") and first.endswith(("章", "节", "条"))
        )
        if is_title:
            section = first.lstrip("# ")
        blocks.append(
            _block(
                path,
                document_id,
                "title" if is_title else "text",
                paragraph,
                index,
                section_title=section,
            )
        )
    return blocks


def _parse_csv(path: Path, document_id: str) -> list[DocumentBlock]:
    with path.open("r", encoding="utf-8-sig", errors="ignore", newline="") as stream:
        rows = list(csv.reader(stream))
    return [
        _block(
            path,
            document_id,
            "table",
            _markdown_table(rows),
            0,
            metadata={"rows": len(rows), "columns": max((len(row) for row in rows), default=0)},
        )
    ]


def _pdf_text_regions(page: Any) -> list[tuple[str, list[float]]]:
    words = sorted(
        page.extract_words(use_text_flow=True, keep_blank_chars=False) or [],
        key=lambda item: (float(item.get("top", 0)), float(item.get("x0", 0))),
    )
    lines: list[dict[str, Any]] = []
    for word in words:
        top = float(word.get("top", 0))
        if not lines or abs(top - lines[-1]["top"]) > 3:
            lines.append(
                {
                    "top": top,
                    "bottom": float(word.get("bottom", top)),
                    "x0": float(word.get("x0", 0)),
                    "x1": float(word.get("x1", 0)),
                    "words": [str(word.get("text", ""))],
                }
            )
        else:
            line = lines[-1]
            line["bottom"] = max(line["bottom"], float(word.get("bottom", top)))
            line["x0"] = min(line["x0"], float(word.get("x0", 0)))
            line["x1"] = max(line["x1"], float(word.get("x1", 0)))
            line["words"].append(str(word.get("text", "")))
    regions: list[list[dict[str, Any]]] = []
    for line in lines:
        if not regions:
            regions.append([line])
            continue
        previous = regions[-1][-1]
        gap = line["top"] - previous["bottom"]
        line_height = max(6.0, previous["bottom"] - previous["top"])
        if gap > line_height * 1.4:
            regions.append([line])
        else:
            regions[-1].append(line)
    result: list[tuple[str, list[float]]] = []
    for region in regions:
        content = "\n".join(" ".join(line["words"]).strip() for line in region).strip()
        if content:
            result.append(
                (
                    content,
                    [
                        min(line["x0"] for line in region),
                        min(line["top"] for line in region),
                        max(line["x1"] for line in region),
                        max(line["bottom"] for line in region),
                    ],
                )
            )
    return result


def _parse_pdf(
    path: Path, document_id: str, asset_dir: Path | None = None
) -> list[DocumentBlock]:
    try:
        import pdfplumber
    except ImportError as exc:
        raise RuntimeError("PDF 结构化解析需要安装 pdfplumber") from exc

    blocks: list[DocumentBlock] = []
    vision_limit = max(0, int(os.getenv("VISION_MAX_IMAGES_PER_DOCUMENT", "20")))
    vision_calls = 0
    ocr_limit = max(0, int(os.getenv("OCR_MAX_PAGES_PER_DOCUMENT", "200")))
    ocr_calls = 0
    pdfium_document = None
    if asset_dir is not None:
        try:
            import pypdfium2 as pdfium

            pdfium_document = pdfium.PdfDocument(str(path))
            asset_dir.mkdir(parents=True, exist_ok=True)
        except ImportError:
            pdfium_document = None
    with pdfplumber.open(path) as pdf:
        for page_number, page in enumerate(pdf.pages, 1):
            block_index = 0
            rendered_page = None
            text_regions = _pdf_text_regions(page)
            page_context = "\n".join(content for content, _ in text_regions)[:1000]
            for paragraph, text_bbox in text_regions:
                blocks.append(
                    _block(
                        path,
                        document_id,
                        "text",
                        paragraph,
                        block_index,
                        page=page_number,
                        bbox=text_bbox,
                    )
                )
                block_index += 1
            try:
                tables = page.find_tables()
            except Exception:
                tables = []
            for table in tables:
                content = _markdown_table(table.extract() or [])
                if content:
                    blocks.append(
                        _block(
                            path,
                            document_id,
                            "table",
                            content,
                            block_index,
                            page=page_number,
                            bbox=table.bbox,
                        )
                    )
                    block_index += 1
            for image_index, image in enumerate(page.images):
                bbox = [
                    image.get("x0", 0),
                    image.get("top", 0),
                    image.get("x1", page.width),
                    image.get("bottom", page.height),
                ]
                asset_path: Path | None = None
                description = ""
                vision_metadata: dict[str, Any] = {"vision_status": "not_rendered"}
                if pdfium_document is not None:
                    if rendered_page is None:
                        rendered_page = pdfium_document[page_number - 1].render(scale=2).to_pil()
                    scale_x = rendered_page.width / float(page.width)
                    scale_y = rendered_page.height / float(page.height)
                    crop_box = (
                        max(0, int(float(bbox[0]) * scale_x)),
                        max(0, int(float(bbox[1]) * scale_y)),
                        min(rendered_page.width, int(float(bbox[2]) * scale_x)),
                        min(rendered_page.height, int(float(bbox[3]) * scale_y)),
                    )
                    if crop_box[2] > crop_box[0] and crop_box[3] > crop_box[1]:
                        asset_path = asset_dir / f"page-{page_number}-image-{image_index + 1}.png"
                        rendered_page.crop(crop_box).save(asset_path)
                        if vision_calls < vision_limit:
                            description, vision_metadata = describe_image(asset_path, page_context)
                            vision_calls += int(vision_metadata.get("vision_status") not in {"disabled", "unconfigured"})
                        else:
                            vision_metadata = {"vision_status": "limit_reached"}
                content = description or f"第{page_number}页图片 {image_index + 1}（待视觉描述）"
                blocks.append(
                    _block(
                        path,
                        document_id,
                        "image",
                        content,
                        block_index,
                        page=page_number,
                        bbox=bbox,
                        metadata={
                            "image_index": image_index,
                            "requires_vision": not bool(description),
                            "asset_name": asset_path.name if asset_path else None,
                            **vision_metadata,
                        },
                    )
                )
                block_index += 1
            if not text_regions and not tables:
                page_asset: Path | None = None
                ocr_text = ""
                ocr_metadata: dict[str, Any] = {"ocr_status": "not_rendered"}
                vision_text = ""
                vision_metadata: dict[str, Any] = {"vision_status": "not_rendered"}
                if pdfium_document is not None:
                    if rendered_page is None:
                        rendered_page = pdfium_document[page_number - 1].render(scale=2).to_pil()
                    page_asset = asset_dir / f"page-{page_number}.png"
                    rendered_page.save(page_asset)
                    if ocr_calls < ocr_limit:
                        ocr_text, ocr_metadata = run_ocr(page_asset)
                        ocr_calls += 1
                    else:
                        ocr_metadata = {"ocr_status": "limit_reached"}
                    if vision_calls < vision_limit:
                        vision_text, vision_metadata = describe_image(page_asset)
                        vision_calls += int(vision_metadata.get("vision_status") not in {"disabled", "unconfigured"})
                    else:
                        vision_metadata = {"vision_status": "limit_reached"}
                blocks.append(
                    _block(
                        path,
                        document_id,
                        "page_image",
                        f"第{page_number}页扫描页面（待 OCR）",
                        block_index,
                        page=page_number,
                        bbox=[0, 0, page.width, page.height],
                        metadata={
                            "requires_ocr": not bool(ocr_text),
                            "requires_vision": not bool(vision_text),
                            "asset_name": page_asset.name if page_asset else None,
                            **ocr_metadata,
                            **vision_metadata,
                        },
                    )
                )
                block_index += 1
                if ocr_text:
                    blocks.append(
                        _block(
                            path,
                            document_id,
                            "ocr_text",
                            ocr_text,
                            block_index,
                            page=page_number,
                            bbox=[0, 0, page.width, page.height],
                            metadata={"asset_name": page_asset.name, **ocr_metadata},
                        )
                    )
                    block_index += 1
                if vision_text:
                    blocks.append(
                        _block(
                            path,
                            document_id,
                            "image_description",
                            vision_text,
                            block_index,
                            page=page_number,
                            bbox=[0, 0, page.width, page.height],
                            metadata={"asset_name": page_asset.name, **vision_metadata},
                        )
                    )
    if pdfium_document is not None:
        pdfium_document.close()
    return blocks


def _parse_docx(
    path: Path, document_id: str, asset_dir: Path | None = None
) -> list[DocumentBlock]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("DOCX 结构化解析需要安装 python-docx") from exc
    document = Document(path)
    blocks: list[DocumentBlock] = []
    section: str | None = None
    index = 0
    for paragraph in document.paragraphs:
        content = paragraph.text.strip()
        if not content:
            continue
        is_title = (paragraph.style and paragraph.style.name.startswith("Heading"))
        if is_title:
            section = content
        blocks.append(
            _block(
                path,
                document_id,
                "title" if is_title else "text",
                content,
                index,
                section_title=section,
            )
        )
        index += 1
    for table in document.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        blocks.append(
            _block(path, document_id, "table", _markdown_table(rows), index, section_title=section)
        )
        index += 1
    image_parts = [
        relation.target_part
        for relation in document.part.rels.values()
        if "image" in str(getattr(relation, "target_ref", ""))
        and hasattr(relation.target_part, "blob")
    ]
    if asset_dir is not None:
        asset_dir.mkdir(parents=True, exist_ok=True)
    vision_limit = max(0, int(os.getenv("VISION_MAX_IMAGES_PER_DOCUMENT", "20")))
    for image_index, image_part in enumerate(image_parts):
        extension = mimetypes.guess_extension(getattr(image_part, "content_type", "")) or ".png"
        asset_path = asset_dir / f"image-{image_index + 1}{extension}" if asset_dir else None
        if asset_path:
            asset_path.write_bytes(image_part.blob)
        if asset_path and image_index < vision_limit:
            description, vision_metadata = describe_image(asset_path, section or "")
        elif asset_path:
            description, vision_metadata = "", {"vision_status": "limit_reached"}
        else:
            description, vision_metadata = "", {"vision_status": "not_extracted"}
        blocks.append(
            _block(
                path,
                document_id,
                "image",
                description or f"文档内图片 {image_index + 1}（待视觉描述）",
                index,
                section_title=section,
                metadata={
                    "image_index": image_index,
                    "requires_vision": not bool(description),
                    "asset_name": asset_path.name if asset_path else None,
                    **vision_metadata,
                },
            )
        )
        index += 1
    return blocks


def _parse_pptx(
    path: Path, document_id: str, asset_dir: Path | None = None
) -> list[DocumentBlock]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise RuntimeError("PPTX 结构化解析需要安装 python-pptx") from exc
    presentation = Presentation(path)
    blocks: list[DocumentBlock] = []
    if asset_dir is not None:
        asset_dir.mkdir(parents=True, exist_ok=True)
    vision_limit = max(0, int(os.getenv("VISION_MAX_IMAGES_PER_DOCUMENT", "20")))
    vision_calls = 0
    for slide_number, slide in enumerate(presentation.slides, 1):
        for shape_index, shape in enumerate(slide.shapes):
            bbox = [shape.left, shape.top, shape.left + shape.width, shape.top + shape.height]
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                blocks.append(
                    _block(path, document_id, "table", _markdown_table(rows), shape_index, page=slide_number, bbox=bbox)
                )
            elif getattr(shape, "has_text_frame", False) and shape.text.strip():
                blocks.append(
                    _block(path, document_id, "text", shape.text, shape_index, page=slide_number, bbox=bbox)
                )
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                extension = str(getattr(shape.image, "ext", "png"))
                asset_path = (
                    asset_dir / f"slide-{slide_number}-image-{shape_index + 1}.{extension}"
                    if asset_dir
                    else None
                )
                if asset_path:
                    asset_path.write_bytes(shape.image.blob)
                if asset_path and vision_calls < vision_limit:
                    description, vision_metadata = describe_image(asset_path)
                    vision_calls += 1
                elif asset_path:
                    description, vision_metadata = "", {"vision_status": "limit_reached"}
                else:
                    description, vision_metadata = "", {"vision_status": "not_extracted"}
                blocks.append(
                    _block(
                        path,
                        document_id,
                        "image",
                        description or f"第{slide_number}页图片（待视觉描述）",
                        shape_index,
                        page=slide_number,
                        bbox=bbox,
                        metadata={
                            "requires_vision": not bool(description),
                            "asset_name": asset_path.name if asset_path else None,
                            **vision_metadata,
                        },
                    )
                )
    return blocks


def parse_document(
    path: Path, asset_root: Path | None = None
) -> tuple[str, list[DocumentBlock]]:
    document_hash = _file_hash(path)
    document_id = document_hash[:24]
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        blocks = _parse_text(path, document_id)
    elif suffix == ".csv":
        blocks = _parse_csv(path, document_id)
    elif suffix == ".pdf":
        blocks = _parse_pdf(
            path,
            document_id,
            asset_root / document_id if asset_root is not None else None,
        )
    elif suffix == ".docx":
        blocks = _parse_docx(
            path,
            document_id,
            asset_root / document_id if asset_root is not None else None,
        )
    elif suffix == ".pptx":
        blocks = _parse_pptx(
            path,
            document_id,
            asset_root / document_id if asset_root is not None else None,
        )
    else:
        raise ValueError(f"不支持的文档类型：{suffix}")
    return document_hash, blocks


def persist_blocks(output_dir: Path, document_id: str, blocks: list[DocumentBlock]) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    target = output_dir / f"{document_id}.jsonl"
    temp = output_dir / f".{document_id}.{uuid.uuid4().hex}.tmp"
    content = "\n".join(json.dumps(block.to_dict(), ensure_ascii=False) for block in blocks)
    temp.write_text(content + ("\n" if content else ""), encoding="utf-8")
    os.replace(temp, target)
    return target


def parser_capabilities() -> dict[str, Any]:
    return {
        "txt": True,
        "markdown": True,
        "csv": True,
        "pdf_layout": importlib.util.find_spec("pdfplumber") is not None,
        "docx": importlib.util.find_spec("docx") is not None,
        "pptx": importlib.util.find_spec("pptx") is not None,
        "pdf_rendering": importlib.util.find_spec("pypdfium2") is not None,
        **enrichment_capabilities(),
        "notes": {
            "ocr": "扫描页在 RapidOCR 可用时自动 OCR，否则保留 requires_ocr。",
            "vision_description": "仅在显式启用并配置 VISION_API_URL 后调用外部视觉模型。",
        },
    }
